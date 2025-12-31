#!/usr/bin/env python3
"""
AIWardrobe Investor Pitch Deck Generator
Creates a professional YC-style PowerPoint presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
from pptx.dml.color import RGBColor
import os

# Color Palette (Premium Dark Theme)
COLORS = {
    'bg_dark': RGBColor(10, 10, 15),
    'bg_card': RGBColor(20, 20, 32),
    'primary': RGBColor(99, 102, 241),
    'primary_light': RGBColor(129, 140, 248),
    'accent': RGBColor(236, 72, 153),
    'success': RGBColor(16, 185, 129),
    'warning': RGBColor(245, 158, 11),
    'error': RGBColor(239, 68, 68),
    'white': RGBColor(255, 255, 255),
    'text_secondary': RGBColor(180, 180, 200),
    'text_muted': RGBColor(130, 130, 150),
}

def create_slide_background(slide, color=COLORS['bg_dark']):
    """Set slide background to dark color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_text(slide, text, left, top, width, height, font_size=44, bold=True, color=COLORS['white']):
    """Add a title text box"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Inter"
    return txBox

def add_body_text(slide, text, left, top, width, height, font_size=18, color=COLORS['text_secondary'], alignment=PP_ALIGN.LEFT):
    """Add body text"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = "Inter"
    p.alignment = alignment
    return txBox

def add_bullet_points(slide, points, left, top, width, height, font_size=16, color=COLORS['text_secondary']):
    """Add bullet point list"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Inter"
        p.space_after = Pt(8)
    return txBox

def add_card(slide, left, top, width, height, fill_color=COLORS['bg_card']):
    """Add a rounded rectangle card"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_slide_number(slide, number, total=10):
    """Add slide number indicator"""
    add_body_text(
        slide, 
        f"{number:02d} / {total:02d}",
        Inches(0.5), Inches(7), Inches(1), Inches(0.3),
        font_size=10, color=COLORS['text_muted']
    )

def create_pitch_deck():
    """Create the complete pitch deck"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 widescreen
    prs.slide_height = Inches(7.5)
    
    # Use blank layout
    blank_layout = prs.slide_layouts[6]
    
    # ========================================
    # SLIDE 1: Title
    # ========================================
    slide1 = prs.slides.add_slide(blank_layout)
    create_slide_background(slide1)
    
    # Logo emoji
    add_title_text(slide1, "👗", Inches(6), Inches(1.5), Inches(2), Inches(1), font_size=72, bold=False)
    
    # Company name
    add_title_text(slide1, "AIWardrobe", Inches(2), Inches(2.5), Inches(9.33), Inches(1), 
                   font_size=60, bold=True, color=COLORS['primary_light'])
    
    # One-liner
    add_body_text(
        slide1,
        "We use AI to help people get dressed faster by scanning their closet with video and generating personalized outfit recommendations.",
        Inches(2), Inches(3.8), Inches(9.33), Inches(1.2),
        font_size=22, color=COLORS['text_secondary'], alignment=PP_ALIGN.CENTER
    )
    
    # Tags
    add_body_text(
        slide1,
        "📱 Mobile App  •  🤖 AI-Powered  •  🎥 Video Scanning",
        Inches(2), Inches(5.2), Inches(9.33), Inches(0.5),
        font_size=16, color=COLORS['text_muted'], alignment=PP_ALIGN.CENTER
    )
    
    # ========================================
    # SLIDE 2: The Problem
    # ========================================
    slide2 = prs.slides.add_slide(blank_layout)
    create_slide_background(slide2)
    add_slide_number(slide2, 2)
    
    # Header
    add_body_text(slide2, "02", Inches(0.5), Inches(0.5), Inches(0.8), Inches(0.4), 
                  font_size=14, color=COLORS['primary_light'])
    add_title_text(slide2, "The Problem", Inches(1.3), Inches(0.4), Inches(5), Inches(0.6), 
                   font_size=36, color=COLORS['primary_light'])
    
    # Main problem statement
    add_card(slide2, Inches(0.5), Inches(1.3), Inches(12.33), Inches(1.2))
    add_body_text(
        slide2,
        "⏰ Every morning, 1.4 billion people waste 15+ minutes deciding what to wear",
        Inches(0.7), Inches(1.5), Inches(11.93), Inches(0.8),
        font_size=22, color=COLORS['white']
    )
    
    # Pain points
    pain_points = [
        ("👔 Forgotten Clothes", "Average person only wears 20% of their wardrobe.\n$12B worth of clothes sit unused in US closets."),
        ("📸 No Digital Wardrobe", "Manually photographing each item takes 2+ hours.\nNobody does it."),
        ("🤔 Decision Fatigue", '"What goes with what?" — No easy way to see\noutfit combinations.'),
    ]
    
    for i, (title, desc) in enumerate(pain_points):
        x = Inches(0.5 + i * 4.1)
        add_card(slide2, x, Inches(2.8), Inches(3.9), Inches(2.2))
        add_body_text(slide2, title, x + Inches(0.2), Inches(3), Inches(3.5), Inches(0.4), 
                      font_size=18, color=COLORS['white'])
        add_body_text(slide2, desc, x + Inches(0.2), Inches(3.5), Inches(3.5), Inches(1.5), 
                      font_size=13, color=COLORS['text_secondary'])
    
    # Current solutions (bad)
    add_body_text(slide2, "How people solve this today:", Inches(0.5), Inches(5.3), Inches(4), Inches(0.4),
                  font_size=14, color=COLORS['text_muted'])
    add_body_text(slide2, "❌ Buy more clothes    ❌ Ignore most of closet    ❌ Manual spreadsheets", 
                  Inches(0.5), Inches(5.7), Inches(10), Inches(0.4),
                  font_size=16, color=COLORS['error'])
    
    # ========================================
    # SLIDE 3: The Solution
    # ========================================
    slide3 = prs.slides.add_slide(blank_layout)
    create_slide_background(slide3)
    add_slide_number(slide3, 3)
    
    add_body_text(slide3, "03", Inches(0.5), Inches(0.5), Inches(0.8), Inches(0.4), 
                  font_size=14, color=COLORS['primary_light'])
    add_title_text(slide3, "The Solution", Inches(1.3), Inches(0.4), Inches(5), Inches(0.6), 
                   font_size=36, color=COLORS['primary_light'])
    
    # Magic benefit
    add_card(slide3, Inches(0.5), Inches(1.3), Inches(12.33), Inches(1))
    add_body_text(
        slide3,
        "✨ Scan your entire wardrobe in 60 seconds with just a video",
        Inches(0.7), Inches(1.5), Inches(11.93), Inches(0.6),
        font_size=24, color=COLORS['white'], alignment=PP_ALIGN.CENTER
    )
    
    # 3-step process
    steps = [
        ("1", "🎥", "Record", "10-second video of your closet"),
        ("2", "🤖", "AI Detects", "Each item auto-categorized"),
        ("3", "👗", "Get Outfits", "Personalized daily suggestions"),
    ]
    
    for i, (num, emoji, title, desc) in enumerate(steps):
        x = Inches(1 + i * 3.8)
        add_card(slide3, x, Inches(2.6), Inches(3.2), Inches(2))
        add_body_text(slide3, emoji, x + Inches(1.1), Inches(2.8), Inches(1), Inches(0.8), 
                      font_size=40, color=COLORS['white'])
        add_body_text(slide3, title, x + Inches(0.2), Inches(3.6), Inches(2.8), Inches(0.4), 
                      font_size=16, color=COLORS['white'], alignment=PP_ALIGN.CENTER)
        add_body_text(slide3, desc, x + Inches(0.2), Inches(4), Inches(2.8), Inches(0.5), 
                      font_size=12, color=COLORS['text_muted'], alignment=PP_ALIGN.CENTER)
    
    # Features grid
    features = [
        ("🌤️", "Weather-Smart", "Outfits based on today's weather"),
        ("✈️", "Trip Planner", "Pack perfectly for any trip"),
        ("👤", "Virtual Try-On", "See outfits on yourself with AI"),
        ("💬", "AI Stylist", "Chat with your personal fashion AI"),
    ]
    
    for i, (emoji, title, desc) in enumerate(features):
        x = Inches(0.5 + i * 3.2)
        add_card(slide3, x, Inches(5), Inches(2.8), Inches(1.4))
        add_body_text(slide3, emoji, x + Inches(1), Inches(5.1), Inches(1), Inches(0.5), font_size=20)
        add_body_text(slide3, title, x + Inches(0.15), Inches(5.5), Inches(2.5), Inches(0.3), 
                      font_size=13, color=COLORS['white'], alignment=PP_ALIGN.CENTER)
        add_body_text(slide3, desc, x + Inches(0.15), Inches(5.85), Inches(2.5), Inches(0.4), 
                      font_size=10, color=COLORS['text_muted'], alignment=PP_ALIGN.CENTER)
    
    # ========================================
    # SLIDE 4: Traction
    # ========================================
    slide4 = prs.slides.add_slide(blank_layout)
    create_slide_background(slide4)
    add_slide_number(slide4, 4)
    
    add_body_text(slide4, "04", Inches(0.5), Inches(0.5), Inches(0.8), Inches(0.4), 
                  font_size=14, color=COLORS['primary_light'])
    add_title_text(slide4, "Traction", Inches(1.3), Inches(0.4), Inches(3), Inches(0.6), 
                   font_size=36, color=COLORS['primary_light'])
    add_body_text(slide4, "The Most Important Slide", Inches(4.5), Inches(0.5), Inches(4), Inches(0.4),
                  font_size=14, color=COLORS['accent'])
    
    # Key metrics
    metrics = [
        ("Built in", "8 weeks", "from idea to working product"),
        ("📱 18+", "", "Screens Built"),
        ("🤖 5", "", "AI Models Integrated"),
        ("🌍 3", "", "Languages Supported"),
    ]
    
    for i, (top, big, bottom) in enumerate(metrics):
        x = Inches(0.5 + i * 3.2)
        add_card(slide4, x, Inches(1.4), Inches(2.9), Inches(1.8))
        add_body_text(slide4, top, x + Inches(0.15), Inches(1.6), Inches(2.6), Inches(0.4),
                      font_size=14, color=COLORS['text_muted'], alignment=PP_ALIGN.CENTER)
        if big:
            add_body_text(slide4, big, x + Inches(0.15), Inches(2), Inches(2.6), Inches(0.6),
                          font_size=36, color=COLORS['primary_light'], alignment=PP_ALIGN.CENTER)
        add_body_text(slide4, bottom, x + Inches(0.15), Inches(2.7), Inches(2.6), Inches(0.4),
                      font_size=12, color=COLORS['text_secondary'], alignment=PP_ALIGN.CENTER)
    
    # What we've built
    add_card(slide4, Inches(0.5), Inches(3.6), Inches(12.33), Inches(3))
    add_body_text(slide4, "What We've Built", Inches(0.7), Inches(3.8), Inches(4), Inches(0.4),
                  font_size=18, color=COLORS['white'])
    
    built_items = [
        "✅ Full mobile app (iOS + Android)",
        "✅ Video-to-wardrobe AI pipeline",
        "✅ Virtual try-on integration",
        "✅ AI chat assistant",
        "✅ Weather-based suggestions",
        "✅ Trip packing planner",
    ]
    
    for i, item in enumerate(built_items):
        col = i // 3
        row = i % 3
        x = Inches(0.7 + col * 6)
        y = Inches(4.4 + row * 0.65)
        add_body_text(slide4, item, x, y, Inches(5.5), Inches(0.5),
                      font_size=16, color=COLORS['success'])
    
    # ========================================
    # SLIDE 5: Unique Insight
    # ========================================
    slide5 = prs.slides.add_slide(blank_layout)
    create_slide_background(slide5)
    add_slide_number(slide5, 5)
    
    add_body_text(slide5, "05", Inches(0.5), Inches(0.5), Inches(0.8), Inches(0.4), 
                  font_size=14, color=COLORS['primary_light'])
    add_title_text(slide5, "Why Now / Why Us", Inches(1.3), Inches(0.4), Inches(6), Inches(0.6), 
                   font_size=36, color=COLORS['primary_light'])
    
    # Quote card
    add_card(slide5, Inches(0.5), Inches(1.2), Inches(12.33), Inches(1.4))
    add_body_text(
        slide5,
        '"Video + AI finally makes wardrobe digitization effortless. Google Gemini and open-source vision models now process video in seconds, not hours."',
        Inches(0.9), Inches(1.4), Inches(11.53), Inches(1),
        font_size=18, color=COLORS['text_secondary']
    )
    
    # Why Now
    add_card(slide5, Inches(0.5), Inches(2.9), Inches(5.9), Inches(2.8))
    add_body_text(slide5, "⏰ Why Now", Inches(0.7), Inches(3.1), Inches(3), Inches(0.4),
                  font_size=18, color=COLORS['primary_light'])
    why_now = [
        "AI Cost Collapse: Vision AI 100x cheaper than 2020",
        "Gemini/GPT-4V: Video understanding finally works",
        "Mobile GPUs: On-device processing now possible",
        "Gen-Z behavior: OOTD culture demands digital wardrobes",
    ]
    add_bullet_points(slide5, why_now, Inches(0.7), Inches(3.6), Inches(5.5), Inches(2),
                      font_size=13, color=COLORS['text_secondary'])
    
    # Why Us
    add_card(slide5, Inches(6.9), Inches(2.9), Inches(5.9), Inches(2.8))
    add_body_text(slide5, "🎯 Our Unfair Advantage", Inches(7.1), Inches(3.1), Inches(4), Inches(0.4),
                  font_size=18, color=COLORS['primary_light'])
    why_us = [
        "Technical: Video-first approach (competitors use photo-by-photo)",
        "Speed: 60 seconds vs 2+ hours to digitize",
        "Stack: Gemini + SAM2 + FashionCLIP pipeline",
        "Market: Multilingual (EN/RU/UZ) for underserved markets",
    ]
    add_bullet_points(slide5, why_us, Inches(7.1), Inches(3.6), Inches(5.5), Inches(2),
                      font_size=13, color=COLORS['text_secondary'])
    
    # Tech moat
    add_body_text(slide5, "Technical Moat:", Inches(0.5), Inches(6), Inches(2), Inches(0.3),
                  font_size=14, color=COLORS['text_muted'])
    add_body_text(slide5, "SAM2 Segmentation  •  FashionCLIP  •  Grounding DINO  •  Gemini Vision  •  Real-time Processing", 
                  Inches(2.5), Inches(6), Inches(10), Inches(0.3),
                  font_size=14, color=COLORS['primary_light'])
    
    # ========================================
    # SLIDE 6: Market Size
    # ========================================
    slide6 = prs.slides.add_slide(blank_layout)
    create_slide_background(slide6)
    add_slide_number(slide6, 6)
    
    add_body_text(slide6, "06", Inches(0.5), Inches(0.5), Inches(0.8), Inches(0.4), 
                  font_size=14, color=COLORS['primary_light'])
    add_title_text(slide6, "Market Size", Inches(1.3), Inches(0.4), Inches(5), Inches(0.6), 
                   font_size=36, color=COLORS['primary_light'])
    
    # Market pyramid
    markets = [
        (Inches(0.5), Inches(12.33), "$52B", "TAM", "Global Fashion Tech & Wardrobe Management"),
        (Inches(1.5), Inches(10.33), "$4.2B", "SAM", "Digital Wardrobe & Outfit Planning Apps"),
        (Inches(2.5), Inches(8.33), "$120M", "SOM", "AI-First Wardrobe Apps in Year 3"),
    ]
    
    for i, (left_offset, width, value, label, desc) in enumerate(markets):
        y = Inches(1.3 + i * 1.2)
        add_card(slide6, left_offset, y, width, Inches(1))
        add_body_text(slide6, value, left_offset + Inches(0.3), y + Inches(0.15), Inches(2), Inches(0.5),
                      font_size=28, color=COLORS['primary_light'])
        add_body_text(slide6, label, left_offset + Inches(2.5), y + Inches(0.15), Inches(1.5), Inches(0.4),
                      font_size=12, color=COLORS['text_muted'])
        add_body_text(slide6, desc, left_offset + Inches(2.5), y + Inches(0.5), Inches(6), Inches(0.4),
                      font_size=14, color=COLORS['text_secondary'])
    
    # Bottom-up calculation
    add_card(slide6, Inches(0.5), Inches(4.8), Inches(12.33), Inches(2.2))
    add_body_text(slide6, "📊 Bottom-Up Calculation", Inches(0.7), Inches(5), Inches(4), Inches(0.4),
                  font_size=16, color=COLORS['white'])
    
    calc_rows = [
        ("Fashion-conscious smartphone users (18-45)", "400M globally"),
        ("× Would pay for wardrobe management (3%)", "12M users"),
        ("× Average annual subscription ($10/mo)", "$120/year"),
        ("= Addressable Market", "$1.44B/year"),
    ]
    
    for i, (label, value) in enumerate(calc_rows):
        y = Inches(5.5 + i * 0.35)
        color = COLORS['primary_light'] if i == 3 else COLORS['text_secondary']
        add_body_text(slide6, label, Inches(0.7), y, Inches(7), Inches(0.3), font_size=13, color=color)
        add_body_text(slide6, value, Inches(10), y, Inches(2.5), Inches(0.3), font_size=13, color=color)
    
    # ========================================
    # SLIDE 7: Competition
    # ========================================
    slide7 = prs.slides.add_slide(blank_layout)
    create_slide_background(slide7)
    add_slide_number(slide7, 7)
    
    add_body_text(slide7, "07", Inches(0.5), Inches(0.5), Inches(0.8), Inches(0.4), 
                  font_size=14, color=COLORS['primary_light'])
    add_title_text(slide7, "Competition", Inches(1.3), Inches(0.4), Inches(5), Inches(0.6), 
                   font_size=36, color=COLORS['primary_light'])
    
    # Comparison table
    add_card(slide7, Inches(0.5), Inches(1.2), Inches(12.33), Inches(3.2))
    
    # Table headers
    headers = ["Feature", "Stylebook", "Cladwell", "AIWardrobe"]
    for i, header in enumerate(headers):
        x = Inches(0.7 + i * 3)
        color = COLORS['primary_light'] if i == 3 else COLORS['text_muted']
        add_body_text(slide7, header, x, Inches(1.4), Inches(2.8), Inches(0.4), font_size=14, color=color)
    
    # Table rows
    rows = [
        ("Video Scanning", "❌", "❌", "✅"),
        ("AI Detection", "❌", "Basic", "✅ Advanced"),
        ("Virtual Try-On", "❌", "❌", "✅"),
        ("Setup Time", "2+ hours", "1+ hour", "60 seconds"),
    ]
    
    for i, row in enumerate(rows):
        y = Inches(1.9 + i * 0.6)
        for j, cell in enumerate(row):
            x = Inches(0.7 + j * 3)
            color = COLORS['success'] if j == 3 and "✅" in cell else (COLORS['error'] if "❌" in cell else COLORS['text_secondary'])
            if j == 3:
                color = COLORS['success']
            add_body_text(slide7, cell, x, y, Inches(2.8), Inches(0.4), font_size=14, color=color)
    
    # Key differentiators
    add_body_text(slide7, "🏆 Key Differentiator: We're the only app that uses VIDEO scanning — 60x faster setup.", 
                  Inches(0.5), Inches(4.7), Inches(12), Inches(0.5),
                  font_size=18, color=COLORS['primary_light'])
    
    # ========================================
    # SLIDE 8: Business Model
    # ========================================
    slide8 = prs.slides.add_slide(blank_layout)
    create_slide_background(slide8)
    add_slide_number(slide8, 8)
    
    add_body_text(slide8, "08", Inches(0.5), Inches(0.5), Inches(0.8), Inches(0.4), 
                  font_size=14, color=COLORS['primary_light'])
    add_title_text(slide8, "Business Model", Inches(1.3), Inches(0.4), Inches(5), Inches(0.6), 
                   font_size=36, color=COLORS['primary_light'])
    
    # Pricing tiers
    tiers = [
        ("Free", "$0", ["3 wardrobe scans", "Basic suggestions", "Weather integration"], "Acquisition"),
        ("Premium", "$9.99/mo", ["Unlimited scans", "Virtual try-on", "AI stylist chat", "Trip planner", "Analytics"], "Core Revenue"),
        ("Annual", "$79.99/yr", ["All Premium features", "33% savings", "Priority support"], "Retention"),
    ]
    
    for i, (name, price, features, purpose) in enumerate(tiers):
        x = Inches(0.5 + i * 4.2)
        add_card(slide8, x, Inches(1.2), Inches(3.8), Inches(3.8))
        
        add_body_text(slide8, name, x + Inches(0.2), Inches(1.4), Inches(3.4), Inches(0.4),
                      font_size=18, color=COLORS['text_muted'], alignment=PP_ALIGN.CENTER)
        add_body_text(slide8, price, x + Inches(0.2), Inches(1.9), Inches(3.4), Inches(0.6),
                      font_size=30, color=COLORS['white'], alignment=PP_ALIGN.CENTER)
        
        for j, feature in enumerate(features):
            add_body_text(slide8, f"✓ {feature}", x + Inches(0.3), Inches(2.6 + j * 0.4), Inches(3.2), Inches(0.35),
                          font_size=12, color=COLORS['success'])
        
        add_body_text(slide8, purpose, x + Inches(0.2), Inches(4.6), Inches(3.4), Inches(0.3),
                      font_size=11, color=COLORS['text_muted'], alignment=PP_ALIGN.CENTER)
    
    # Unit economics
    add_body_text(slide8, "Unit Economics Target:", Inches(0.5), Inches(5.4), Inches(3), Inches(0.4),
                  font_size=14, color=COLORS['text_muted'])
    add_body_text(slide8, "LTV: $180  •  CAC: $15  •  LTV/CAC: 12x", Inches(3.5), Inches(5.4), Inches(8), Inches(0.4),
                  font_size=18, color=COLORS['success'])
    
    # Revenue mix
    add_body_text(slide8, "Revenue Mix:  70% Subscriptions  •  20% Affiliate  •  10% Brand Partnerships", 
                  Inches(0.5), Inches(6), Inches(12), Inches(0.4),
                  font_size=14, color=COLORS['text_secondary'])
    
    # ========================================
    # SLIDE 9: Team
    # ========================================
    slide9 = prs.slides.add_slide(blank_layout)
    create_slide_background(slide9)
    add_slide_number(slide9, 9)
    
    add_body_text(slide9, "09", Inches(0.5), Inches(0.5), Inches(0.8), Inches(0.4), 
                  font_size=14, color=COLORS['primary_light'])
    add_title_text(slide9, "The Team", Inches(1.3), Inches(0.4), Inches(5), Inches(0.6), 
                   font_size=36, color=COLORS['primary_light'])
    
    # Founder card
    add_card(slide9, Inches(3.5), Inches(1.3), Inches(6.33), Inches(3.5))
    
    # Avatar placeholder
    avatar_shape = slide9.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.9), Inches(1.5), Inches(1.5), Inches(1.5))
    avatar_shape.fill.solid()
    avatar_shape.fill.fore_color.rgb = COLORS['primary']
    avatar_shape.line.fill.background()
    add_body_text(slide9, "ZV", Inches(6.1), Inches(1.9), Inches(1.1), Inches(0.7), font_size=28, color=COLORS['white'])
    
    add_body_text(slide9, "Zohid Vohidjonov", Inches(3.7), Inches(3.1), Inches(5.93), Inches(0.5),
                  font_size=24, color=COLORS['white'], alignment=PP_ALIGN.CENTER)
    add_body_text(slide9, "Founder & Full-Stack Developer", Inches(3.7), Inches(3.6), Inches(5.93), Inches(0.4),
                  font_size=14, color=COLORS['primary_light'], alignment=PP_ALIGN.CENTER)
    
    credentials = [
        "🎓 Self-taught developer",
        "💻 React Native + Express.js",
        "🤖 AI/ML Integration",
        "🚀 Built full product solo in 8 weeks"
    ]
    add_body_text(slide9, "  •  ".join(credentials), Inches(3.7), Inches(4.1), Inches(5.93), Inches(0.4),
                  font_size=11, color=COLORS['text_secondary'], alignment=PP_ALIGN.CENTER)
    
    # Hiring roadmap
    add_card(slide9, Inches(0.5), Inches(5.2), Inches(12.33), Inches(1.5))
    add_body_text(slide9, "Hiring Roadmap with Funding", Inches(0.7), Inches(5.4), Inches(5), Inches(0.4),
                  font_size=16, color=COLORS['white'])
    
    hires = [
        ("Q1", "AI/ML Engineer"),
        ("Q1", "Growth Marketer"),
        ("Q2", "iOS Developer"),
        ("Q2", "Product Designer"),
    ]
    
    for i, (quarter, role) in enumerate(hires):
        x = Inches(0.7 + i * 3)
        add_body_text(slide9, quarter, x, Inches(5.95), Inches(0.5), Inches(0.3), font_size=12, color=COLORS['primary'])
        add_body_text(slide9, role, x + Inches(0.6), Inches(5.95), Inches(2.2), Inches(0.3), font_size=13, color=COLORS['text_secondary'])
    
    # ========================================
    # SLIDE 10: The Ask
    # ========================================
    slide10 = prs.slides.add_slide(blank_layout)
    create_slide_background(slide10)
    add_slide_number(slide10, 10)
    
    add_body_text(slide10, "10", Inches(0.5), Inches(0.5), Inches(0.8), Inches(0.4), 
                  font_size=14, color=COLORS['primary_light'])
    add_title_text(slide10, "The Ask", Inches(1.3), Inches(0.4), Inches(5), Inches(0.6), 
                   font_size=36, color=COLORS['primary_light'])
    
    # Big ask number
    add_title_text(slide10, "$500K", Inches(2), Inches(1.2), Inches(9.33), Inches(1),
                   font_size=80, color=COLORS['primary_light'])
    add_body_text(slide10, "Pre-Seed Round", Inches(2), Inches(2.2), Inches(9.33), Inches(0.4),
                  font_size=18, color=COLORS['text_muted'], alignment=PP_ALIGN.CENTER)
    
    # Use of funds
    add_body_text(slide10, "Use of Funds", Inches(0.5), Inches(2.9), Inches(3), Inches(0.4),
                  font_size=16, color=COLORS['white'])
    
    funds = [
        ("Engineering", "$225K (45%)", "AI/ML + Mobile team"),
        ("Growth & Marketing", "$150K (30%)", "User acquisition"),
        ("Infrastructure", "$75K (15%)", "Servers, AI APIs"),
        ("Operations", "$50K (10%)", "Legal, admin"),
    ]
    
    for i, (category, amount, desc) in enumerate(funds):
        y = Inches(3.4 + i * 0.5)
        add_body_text(slide10, category, Inches(0.5), y, Inches(3), Inches(0.4), font_size=13, color=COLORS['text_secondary'])
        add_body_text(slide10, amount, Inches(3.8), y, Inches(2), Inches(0.4), font_size=13, color=COLORS['primary_light'])
        add_body_text(slide10, desc, Inches(6), y, Inches(3), Inches(0.4), font_size=13, color=COLORS['text_muted'])
    
    # Milestones
    add_body_text(slide10, "18-Month Milestones", Inches(0.5), Inches(5.5), Inches(4), Inches(0.4),
                  font_size=16, color=COLORS['white'])
    
    milestones = [
        ("👥", "50,000", "Active Users"),
        ("💰", "$30K", "MRR"),
        ("🌍", "10", "Languages"),
        ("🤝", "5", "Brand Partnerships"),
    ]
    
    for i, (emoji, value, label) in enumerate(milestones):
        x = Inches(0.5 + i * 3.2)
        add_body_text(slide10, emoji, x + Inches(0.8), Inches(5.9), Inches(0.6), Inches(0.5), font_size=20)
        add_body_text(slide10, value, x, Inches(6.3), Inches(2.8), Inches(0.5), font_size=24, color=COLORS['primary_light'], alignment=PP_ALIGN.CENTER)
        add_body_text(slide10, label, x, Inches(6.75), Inches(2.8), Inches(0.3), font_size=11, color=COLORS['text_muted'], alignment=PP_ALIGN.CENTER)
    
    # Save the presentation
    output_path = "/Users/zohidvohidjonov/Desktop/AIWardrobe/AIWardrobe_Investor_Pitch_YC.pptx"
    prs.save(output_path)
    print(f"✅ Pitch deck saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    create_pitch_deck()
